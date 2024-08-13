#!/usr/bin/env python3
import os
import sys
import venv
import subprocess

VENV_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), "file_search_venv")
REQUIREMENTS = [
    "PyMuPDF",
    "python-pptx",
    "prompt_toolkit",
    "openpyxl",
    "python-docx",
    "pandas",
    "lxml",
    "beautifulsoup4",
    "markdown",
]

def create_venv_if_needed():
    if not os.path.exists(VENV_DIR):
        print("Creating virtual environment...")
        venv.create(VENV_DIR, with_pip=True)
        
        pip_path = os.path.join(VENV_DIR, "Scripts", "pip.exe") if sys.platform == "win32" else os.path.join(VENV_DIR, "bin", "pip")
        
        print("Installing required packages...")
        subprocess.check_call([pip_path, "install"] + REQUIREMENTS)
        print("Virtual environment setup complete.")
    else:
        print("Using existing virtual environment.")

def run_script():
    python_path = os.path.join(VENV_DIR, "Scripts", "python.exe") if sys.platform == "win32" else os.path.join(VENV_DIR, "bin", "python")
    
    script_path = os.path.abspath(__file__)
    os.execv(python_path, [python_path, script_path, "run_in_venv"])

if __name__ == "__main__":
    if len(sys.argv) > 1 and sys.argv[1] == "run_in_venv":
        import signal
        import multiprocessing
        from functools import partial
        import mmap
        import fitz  # PyMuPDF
        from pptx import Presentation
        from prompt_toolkit import PromptSession
        from prompt_toolkit.completion import PathCompleter
        from prompt_toolkit.history import FileHistory
        import openpyxl
        from docx import Document
        import pandas as pd
        import xml.etree.ElementTree as ET
        from bs4 import BeautifulSoup
        import markdown
        import json

        # Global variables
        current_search_dir = ""
        session = PromptSession(history=FileHistory(os.path.expanduser("~/.file_search_history")))

        def search_pdf(file_path, search_string):
            try:
                with fitz.open(file_path) as doc:
                    for page in doc:
                        if search_string in page.get_text().lower():
                            return True
            except Exception as e:
                print(f"Error processing PDF {file_path}: {e}")
            return False

        def search_pptx(file_path, search_string):
            try:
                prs = Presentation(file_path)
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, 'text'):
                            if search_string in shape.text.lower():
                                return True
            except Exception as e:
                print(f"Error processing PPTX {file_path}: {e}")
            return False

        def search_excel(file_path, search_string):
            try:
                workbook = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
                for sheet in workbook.sheetnames:
                    for row in workbook[sheet].iter_rows():
                        for cell in row:
                            if cell.value and search_string in str(cell.value).lower():
                                return True
            except Exception as e:
                print(f"Error processing Excel {file_path}: {e}")
            return False

        def search_word(file_path, search_string):
            try:
                doc = Document(file_path)
                for para in doc.paragraphs:
                    if search_string in para.text.lower():
                        return True
            except Exception as e:
                print(f"Error processing Word {file_path}: {e}")
            return False

        def search_text(file_path, search_string):
            try:
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as file:
                    return search_string in file.read().lower()
            except Exception as e:
                print(f"Error processing text file {file_path}: {e}")
            return False

        def search_csv(file_path, search_string):
            try:
                df = pd.read_csv(file_path)
                return df.astype(str).apply(lambda x: x.str.contains(search_string, case=False)).any().any()
            except Exception as e:
                print(f"Error processing CSV {file_path}: {e}")
            return False

        def search_xml(file_path, search_string):
            try:
                tree = ET.parse(file_path)
                root = tree.getroot()
                for elem in root.iter():
                    if elem.text and isinstance(elem.text, str) and search_string in elem.text.lower():
                        return True
            except Exception as e:
                print(f"Error processing XML {file_path}: {e}")
            return False

        def search_html(file_path, search_string):
            try:
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as file:
                    soup = BeautifulSoup(file, 'lxml')
                    return search_string in soup.get_text().lower()
            except Exception as e:
                print(f"Error processing HTML {file_path}: {e}")
            return False

        def search_markdown(file_path, search_string):
            try:
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as file:
                    html = markdown.markdown(file.read())
                    return search_string in html.lower()
            except Exception as e:
                print(f"Error processing Markdown {file_path}: {e}")
            return False

        def search_json(file_path, search_string):
            try:
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as file:
                    data = json.load(file)
                    return search_string in json.dumps(data).lower()
            except Exception as e:
                print(f"Error processing JSON {file_path}: {e}")
            return False

        def worker(file_path, search_string):
            ext = os.path.splitext(file_path)[1].lower()
            if ext == '.pdf':
                return file_path if search_pdf(file_path, search_string) else None
            elif ext == '.pptx':
                return file_path if search_pptx(file_path, search_string) else None
            elif ext in ('.xlsx', '.xls'):
                return file_path if search_excel(file_path, search_string) else None
            elif ext in ('.docx', '.doc'):
                return file_path if search_word(file_path, search_string) else None
            elif ext in ('.txt', '.rtf', '.odt', '.ods'):
                return file_path if search_text(file_path, search_string) else None
            elif ext == '.csv':
                return file_path if search_csv(file_path, search_string) else None
            elif ext == '.xml':
                return file_path if search_xml(file_path, search_string) else None
            elif ext in ('.html', '.htm'):
                return file_path if search_html(file_path, search_string) else None
            elif ext == '.md':
                return file_path if search_markdown(file_path, search_string) else None
            elif ext == '.json':
                return file_path if search_json(file_path, search_string) else None
            elif os.path.getsize(file_path) > 10 * 1024 * 1024:  # Files larger than 10MB
                return file_path if search_text(file_path, search_string) else None
            return None

        def search_files(directory, search_string):
            search_string = search_string.lower()
            file_list = []
            for root, _, files in os.walk(directory):
                for file in files:
                    ext = os.path.splitext(file)[1].lower()
                    if ext in ('.pdf', '.pptx', '.xlsx', '.xls', '.docx', '.doc', '.txt', '.rtf', '.odt', '.ods', '.csv', '.xml', '.html', '.htm', '.md', '.json'):
                        file_list.append(os.path.join(root, file))
                    elif os.path.getsize(os.path.join(root, file)) > 10 * 1024 * 1024:
                        file_list.append(os.path.join(root, file))

            with multiprocessing.Pool() as pool:
                results = pool.map(partial(worker, search_string=search_string), file_list)
            
            return [r for r in results if r is not None]

        def get_valid_directory():
            while True:
                try:
                    search_dir = session.prompt(
                        "Enter the directory path to search: ",
                        completer=PathCompleter(),
                        complete_while_typing=True
                    ).strip()

                    if os.path.exists(search_dir) and os.path.isdir(search_dir):
                        if os.access(search_dir, os.R_OK):
                            return search_dir
                        else:
                            print("Error: You don't have read permissions for this directory.")
                    else:
                        print("Error: The directory doesn't exist or is not a valid directory.")
                    print("Please try again.")
                except KeyboardInterrupt:
                    print("\nOperation cancelled. Please try again.")

        def change_directory(signum, frame):
            global current_search_dir
            print("\n\nChanging search directory...")
            current_search_dir = get_valid_directory()
            print(f"\nDirectory changed to: {current_search_dir}")
            print("You can now perform searches in the new directory.")
            print("Press Ctrl+C to exit, or Ctrl+Z to change directory again.")

        def main():
            global current_search_dir
            current_search_dir = get_valid_directory()
            print(f"\nDirectory set to: {current_search_dir}")
            print("You can now perform multiple case-insensitive searches in this directory.")
            print("Supported file types: PDF, PPTX, XLSX, XLS, DOCX, DOC, TXT, RTF, ODT, ODS, CSV, XML, HTML, HTM, MD, JSON")
            print("Press Ctrl+C to exit, or Ctrl+Z to change directory.")

            signal.signal(signal.SIGTSTP, change_directory)

            try:
                while True:
                    search_term = session.prompt("\nEnter the search term (or press Ctrl+C to exit): ").strip()
                    
                    print(f"\nSearching for '{search_term}' (case-insensitive) in '{current_search_dir}'...")
                    found_files = search_files(current_search_dir, search_term)

                    if found_files:
                        print(f"\nFiles containing '{search_term}' (case-insensitive):")
                        for file in found_files:
                            print(file)
                    else:
                        print(f"\nNo files containing '{search_term}' (case-insensitive) were found.")
                    
                    print("\n" + "-"*50)

            except KeyboardInterrupt:
                print("\n\nSearch terminated by user. Goodbye!")

        if __name__ == "__main__":
            main()

    else:
        create_venv_if_needed()
        run_script()
